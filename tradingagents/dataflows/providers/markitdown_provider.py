"""MarkItDown ingestion: microsoft/markitdown wrapper with graceful fallback."""

from __future__ import annotations

import csv
import io
import json
import logging
import pathlib
import re
import zipfile
from typing import Any

logger = logging.getLogger(__name__)

try:
    from markitdown import MarkItDown as _MarkItDown  # type: ignore

    _HAS_MARKITDOWN = True
except Exception:  # noqa: BLE001
    _MarkItDown = None  # type: ignore
    _HAS_MARKITDOWN = False

_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff"}
_TABULAR_EXTS = {".xlsx", ".xls"}
_MAX_CSV_ROWS = 200
_MAX_CHARS = 120_000


class MarkItDownProvider:
    """Convert PDF/PPTX/XLSX/HTML/CSV/JSON/ZIP/images to markdown.

    Uses ``markitdown`` when installed, else lightweight fallbacks.
    ``llm_client`` enables image description (markitdown plugin pattern).
    """

    def __init__(self, llm_client: Any | None = None, enable_llm_image: bool = False) -> None:
        self.llm_client = llm_client
        self.enable_llm_image = enable_llm_image
        self._md: Any | None = None
        if _HAS_MARKITDOWN:
            try:
                kw: dict[str, Any] = {}
                if llm_client is not None:
                    kw["llm_client"] = llm_client
                    kw["llm_model"] = getattr(llm_client, "model", None)
                self._md = _MarkItDown(**kw) if kw else _MarkItDown()
            except Exception as exc:  # noqa: BLE001
                logger.warning("MarkItDown init failed, fallback: %s", exc)
        else:
            logger.warning("markitdown not installed; using fallback parsers")

    def convert(self, path: str | bytes, ext: str | None = None) -> str:
        """Convert document to markdown. ``path`` is file path or raw bytes."""
        resolved_ext, data, file_path = self._resolve_input(path, ext)
        if self._md is not None and file_path is not None:
            try:
                result = self._md.convert(str(file_path))
                text = getattr(result, "text_content", str(result))
                if text and text.strip():
                    return text
            except Exception as exc:  # noqa: BLE001
                logger.warning("markitdown convert failed for %s: %s", file_path, exc)
        return self._fallback_convert(resolved_ext, data, file_path)

    def convert_for_llm(self, path: str | bytes | pathlib.Path) -> str:
        """Markdown ready for LLM consumption (preserves tables/headings)."""
        src = str(path) if isinstance(path, pathlib.Path) else (path if isinstance(path, str) else "<bytes>")
        try:
            md = self.convert(path)  # type: ignore[arg-type]
        except Exception as exc:  # noqa: BLE001
            logger.warning("convert_for_llm failed for %s: %s", src, exc)
            return f"# Document: {src}\n\n_CONVERSION_FAILED: {exc}_"
        content = f"# Document: {src}\n\n" + (md or "_EMPTY_DOCUMENT_")
        if len(content) > _MAX_CHARS:
            content = content[:_MAX_CHARS] + f"\n\n_...truncated ({len(content)} chars)_"
        return content

    # -- internals --
    def _resolve_input(self, path: str | bytes, ext: str | None):
        if isinstance(path, bytes):
            e = (ext or "").lower()
            if e and not e.startswith("."):
                e = f".{e}"
            return e or ".txt", path, None
        p = pathlib.Path(path)
        e = (ext or p.suffix or "").lower()
        if e and not e.startswith("."):
            e = f".{e}"
        data: bytes | None = None
        fp: pathlib.Path | None = None
        if p.exists() and p.is_file():
            fp = p
            try:
                data = p.read_bytes()
            except Exception:
                data = None
        else:
            data = path.encode("utf-8", errors="replace")
            if not e:
                e = ".txt"
        return e, data, fp

    def _fallback_convert(self, ext: str, data: bytes | None, fp: pathlib.Path | None) -> str:
        if ext in _IMAGE_EXTS:
            return self._image_fallback(fp, ext)
        if ext == ".pdf":
            return self._pdf_fallback(data, fp)
        if ext in _TABULAR_EXTS:
            return self._excel_fallback(data, fp)
        if ext == ".csv":
            return self._csv_fallback(data, fp)
        if ext == ".json":
            return self._json_fallback(data, fp)
        if ext in {".html", ".htm"}:
            return self._html_fallback(data, fp)
        if ext == ".pptx":
            return self._pptx_fallback(data, fp)
        if ext == ".zip":
            return self._zip_fallback(data, fp)
        return self._text_fallback(data, fp)

    def _read_text(self, data: bytes | None, fp: pathlib.Path | None) -> str:
        if data is not None:
            return data.decode("utf-8", errors="replace")
        if fp and fp.exists():
            return fp.read_text(encoding="utf-8", errors="replace")
        return ""

    def _pdf_fallback(self, data: bytes | None, fp: pathlib.Path | None) -> str:
        try:
            import PyPDF2  # type: ignore

            reader = PyPDF2.PdfReader(io.BytesIO(data)) if data is not None else PyPDF2.PdfReader(str(fp))
            text = "\n\n".join((p.extract_text() or "") for p in reader.pages).strip()
            if text:
                return text
        except Exception as exc:  # noqa: BLE001
            logger.debug("PyPDF2 failed: %s", exc)
        raw = self._read_text(data, fp)
        return raw if raw.strip() else "_EMPTY_PDF_"

    def _excel_fallback(self, data: bytes | None, fp: pathlib.Path | None) -> str:
        try:
            import openpyxl  # type: ignore

            wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True) if data else openpyxl.load_workbook(str(fp), read_only=True, data_only=True)
            parts: list[str] = []
            for ws in wb.worksheets:
                rows = list(ws.iter_rows(values_only=True))
                if not rows:
                    continue
                parts.append(f"## Sheet: {ws.title}")
                hdr = [str(c) if c is not None else "" for c in rows[0]]
                parts.append("| " + " | ".join(hdr) + " |")
                parts.append("| " + " | ".join(["---"] * len(hdr)) + " |")
                for r in rows[1:]:
                    parts.append("| " + " | ".join(str(c) if c is not None else "" for c in r) + " |")
                parts.append("")
            return "\n".join(parts).strip() or "_EMPTY_EXCEL_"
        except Exception as exc:  # noqa: BLE001
            logger.debug("openpyxl failed: %s", exc)
            return self._read_text(data, fp) or "_EXCEL_PARSE_FAILED_"

    def _csv_fallback(self, data: bytes | None, fp: pathlib.Path | None) -> str:
        text = self._read_text(data, fp)
        if not text.strip():
            return "_EMPTY_CSV_"
        rows = list(csv.reader(io.StringIO(text)))[: _MAX_CSV_ROWS + 1]
        if not rows:
            return text
        out = ["| " + " | ".join(rows[0]) + " |", "| " + " | ".join(["---"] * len(rows[0])) + " |"]
        out += ["| " + " | ".join(r) + " |" for r in rows[1:]]
        if len(rows) > _MAX_CSV_ROWS:
            out.append(f"\n_...truncated, {len(rows)} rows_")
        return "\n".join(out)

    def _json_fallback(self, data: bytes | None, fp: pathlib.Path | None) -> str:
        text = self._read_text(data, fp)
        try:
            pretty = json.dumps(json.loads(text), indent=2, ensure_ascii=False)
            return f"```json\n{pretty}\n```"
        except Exception:
            return f"```json\n{text}\n```" if text.strip() else "_EMPTY_JSON_"

    def _html_fallback(self, data: bytes | None, fp: pathlib.Path | None) -> str:
        text = self._read_text(data, fp)
        try:
            from bs4 import BeautifulSoup  # type: ignore

            return BeautifulSoup(text, "html.parser").get_text(separator="\n", strip=True) or text
        except Exception:
            return re.sub(r"\s+", " ", re.sub(r"<[^>]+>", " ", text)).strip() or text

    def _pptx_fallback(self, data: bytes | None, fp: pathlib.Path | None) -> str:
        try:
            import pptx  # type: ignore

            prs = pptx.Presentation(io.BytesIO(data)) if data else pptx.Presentation(str(fp))
            slides = []
            for i, s in enumerate(prs.slides, 1):
                txt = [f"## Slide {i}"] + [sh.text for sh in s.shapes if hasattr(sh, "text") and sh.text]
                slides.append("\n".join(txt))
            return "\n\n".join(slides) if slides else "_EMPTY_PPTX_"
        except Exception as exc:  # noqa: BLE001
            logger.debug("pptx failed: %s", exc)
            return self._read_text(data, fp) or "_PPTX_PARSE_FAILED_"

    def _zip_fallback(self, data: bytes | None, fp: pathlib.Path | None) -> str:
        try:
            buf = io.BytesIO(data) if data is not None else open(str(fp), "rb")  # noqa: SIM115
            close = isinstance(buf, io.BufferedReader)
            with zipfile.ZipFile(buf) as zf:
                entries = zf.namelist()
                parts = [f"# ZIP: {fp or '<bytes>'}", f"Contents: {', '.join(entries[:20])}", ""]
                for name in entries[:10]:
                    try:
                        inner = zf.read(name)
                        ext = pathlib.Path(name).suffix.lower()
                        parts.append(f"## {name}\n" + self._fallback_convert(ext, inner, None))
                    except Exception as exc:  # noqa: BLE001
                        parts.append(f"## {name}\n_FAILED: {exc}_")
                return "\n".join(parts)
        except Exception as exc:  # noqa: BLE001
            return f"_ZIP_PARSE_FAILED: {exc}_"
        finally:
            if "close" in locals() and close:
                try:
                    buf.close()  # type: ignore
                except Exception:
                    pass

    def _image_fallback(self, fp: pathlib.Path | None, ext: str) -> str:
        if self.enable_llm_image and self.llm_client is not None and fp:
            try:
                prompt = f"Describe this financial chart/image: {fp.name}"
                if hasattr(self.llm_client, "invoke"):
                    resp = self.llm_client.invoke(prompt)
                    desc = getattr(resp, "content", str(resp))
                    return f"![{fp.name}]({fp.name})\n\n{desc}"
                if callable(self.llm_client):
                    return f"![{fp.name}]({fp.name})\n\n{self.llm_client(prompt)}"
            except Exception as exc:  # noqa: BLE001
                logger.debug("LLM image failed: %s", exc)
        name = fp.name if fp else f"image{ext}"
        return f"![{name}]({name})\n\n_Image placeholder for {name}; install markitdown or provide llm_client_"

    def _text_fallback(self, data: bytes | None, fp: pathlib.Path | None) -> str:
        return self._read_text(data, fp) or "_EMPTY_DOCUMENT_"
